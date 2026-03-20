# exports.py — Generadores de reportes: Excel, PDF, PowerPoint
from io import BytesIO
from datetime import date, datetime


# ── Colores corporativos (consistentes con style.css) ─────────────
_C = {
    'primary':  '#4F46E5',
    'success':  '#16A34A',
    'warning':  '#D97706',
    'danger':   '#DC2626',
    'info':     '#0284C7',
    'muted':    '#6B7280',
    'bg':       '#EEF2FF',
    'white':    '#FFFFFF',
    'dark':     '#1E293B',
    'row_alt':  '#F9FAFB',
    'border':   '#E5E7EB',
}


def _fmt_fecha(val):
    if val is None:
        return '—'
    if isinstance(val, (date, datetime)):
        return val.strftime('%d/%m/%Y')
    s = str(val)[:10]
    return s if s and s != 'None' else '—'


def _sem_label(sem):
    return {'verde': 'Verde', 'amarillo': 'Amarillo', 'rojo': 'Rojo'}.get(sem, '—')


def _stats_dict(filtered, sem_verde, sem_amarillo, sem_rojo,
                pendientes_frd, pendientes_firma, pendientes_aceptacion, frd_pct, total_frd):
    return {
        'total': len(filtered),
        'sem_verde': sem_verde, 'sem_amarillo': sem_amarillo, 'sem_rojo': sem_rojo,
        'pendientes_frd': pendientes_frd, 'pendientes_firma': pendientes_firma,
        'pendientes_aceptacion': pendientes_aceptacion,
        'frd_pct': frd_pct, 'total_frd': total_frd,
    }


# =============================================================================
#  EXCEL
# =============================================================================
def generate_excel(filtered, division_filtro, sem_verde, sem_amarillo, sem_rojo,
                   pendientes_frd, pendientes_firma, pendientes_aceptacion, frd_pct, total_frd):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, GradientFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    def thin_border():
        s = Side(style='thin', color='D1D5DB')
        return Border(left=s, right=s, top=s, bottom=s)

    def header_cell(ws, row, col, value, bg=_C['primary'], fg=_C['white'], size=9, bold=True):
        c = ws.cell(row=row, column=col, value=value)
        c.font = Font(bold=bold, size=size, color=fg.lstrip('#'))
        c.fill = PatternFill('solid', start_color=bg.lstrip('#'))
        c.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        c.border = thin_border()
        return c

    def data_cell(ws, row, col, value, color=None, bold=False, align='center'):
        c = ws.cell(row=row, column=col, value=value)
        c.font = Font(color=(color or _C['dark']).lstrip('#'), bold=bold, size=9)
        c.alignment = Alignment(horizontal=align, vertical='center', wrap_text=(col in (2, 12)))
        c.border = thin_border()
        return c

    # ── Hoja 1: Resumen ───────────────────────────────────────────
    ws = wb.active
    ws.title = 'Resumen'

    ws['A1'] = 'Gestor de Entregables — Reporte ejecutivo'
    ws['A1'].font = Font(bold=True, size=16, color=_C['primary'].lstrip('#'))
    ws['A2'] = f"División: {division_filtro or 'Todas las divisiones'}"
    ws['A2'].font = Font(size=11, color=_C['muted'].lstrip('#'))
    ws['A3'] = f"Generado: {datetime.now().strftime('%d/%m/%Y  %H:%M')}"
    ws['A3'].font = Font(size=9, color=_C['muted'].lstrip('#'), italic=True)
    ws.row_dimensions[1].height = 26

    ws['A5'] = 'Indicador'
    ws['B5'] = 'Valor'
    header_cell(ws, 5, 1, 'Indicador')
    header_cell(ws, 5, 2, 'Valor')

    kpis = [
        ('Total entregables',          len(filtered),           _C['primary']),
        ('Semáforo  ●  Verde',          sem_verde,               _C['success']),
        ('Semáforo  ●  Amarillo',       sem_amarillo,            _C['warning']),
        ('Semáforo  ●  Rojo',           sem_rojo,                _C['danger']),
        ('FRD — Sin URL',               pendientes_frd,          _C['warning']),
        ('FRD — Sin firma',             pendientes_firma,        _C['info']),
        ('FRD — Sin aceptación',        pendientes_aceptacion,   _C['danger']),
        (f'FRDs completados ({total_frd} total)', f'{frd_pct}%', _C['success']),
    ]
    for i, (label, value, color) in enumerate(kpis, 6):
        ws.cell(row=i, column=1, value=label).border = thin_border()
        c = ws.cell(row=i, column=2, value=value)
        c.font = Font(bold=True, size=11, color=color.lstrip('#'))
        c.alignment = Alignment(horizontal='center', vertical='center')
        c.border = thin_border()
        if i % 2 == 0:
            for col in (1, 2):
                ws.cell(row=i, column=col).fill = PatternFill('solid', start_color='F9FAFB')

    ws.column_dimensions['A'].width = 32
    ws.column_dimensions['B'].width = 16

    # ── Hoja 2: Entregables ───────────────────────────────────────
    ws2 = wb.create_sheet('Entregables')

    hdrs = ['Folio', 'Nombre', 'Estatus', 'Semáforo', 'Avance %',
            'Responsable', 'División', 'Key User',
            'Fecha Inicio', 'Fecha Entrega', 'FRD Aplica', 'Descripción']
    for col, h in enumerate(hdrs, 1):
        header_cell(ws2, 1, col, h)
    ws2.row_dimensions[1].height = 22
    ws2.freeze_panes = 'A2'

    sem_colors = {'verde': _C['success'], 'amarillo': _C['warning'], 'rojo': _C['danger']}

    for r, e in enumerate(filtered, 2):
        sem = e.get('semaforo_cliente', '')
        avance = e.get('porcentaje_avance') or 0
        avance_color = _C['success'] if avance >= 80 else (_C['warning'] if avance >= 40 else _C['danger'])

        row_bg = _C['row_alt'] if r % 2 == 0 else _C['white']
        vals = [
            (e.get('folio') or '', _C['dark'], False, 'center'),
            (e.get('nombre') or '', _C['dark'], False, 'left'),
            (e.get('estatus') or '', _C['dark'], False, 'center'),
            (_sem_label(sem), sem_colors.get(sem, _C['dark']), True, 'center'),
            (avance, avance_color, True, 'center'),
            (e.get('responsable') or '', _C['dark'], False, 'left'),
            (e.get('division') or '', _C['dark'], False, 'center'),
            (e.get('key_user') or '', _C['dark'], False, 'center'),
            (_fmt_fecha(e.get('fecha_inicio')), _C['muted'], False, 'center'),
            (_fmt_fecha(e.get('fecha_entrega')), _C['muted'], False, 'center'),
            ('Sí' if e.get('frd_aplica') else 'No', _C['dark'], False, 'center'),
            (e.get('descripcion') or '', _C['muted'], False, 'left'),
        ]
        for col, (val, color, bold, align) in enumerate(vals, 1):
            c = data_cell(ws2, r, col, val, color, bold, align)
            c.fill = PatternFill('solid', start_color=row_bg.lstrip('#'))

    col_widths = [10, 38, 14, 12, 10, 26, 18, 18, 14, 14, 12, 50]
    for i, w in enumerate(col_widths, 1):
        ws2.column_dimensions[get_column_letter(i)].width = w

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
#  PDF
# =============================================================================
def generate_pdf(filtered, division_filtro, sem_verde, sem_amarillo, sem_rojo,
                 pendientes_frd, pendientes_firma, pendientes_aceptacion, frd_pct, total_frd):
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                    Paragraph, Spacer, HRFlowable)
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT

    def hex2color(h):
        h = h.lstrip('#')
        return colors.HexColor(f'#{h}')

    C = {k: hex2color(v) for k, v in _C.items()}
    C['light_bg'] = hex2color('#EEF2FF')

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=landscape(A4),
        leftMargin=1.5*cm, rightMargin=1.5*cm,
        topMargin=1.5*cm, bottomMargin=1.5*cm,
        title='Reporte de Entregables',
    )

    title_s  = ParagraphStyle('t', fontSize=20, textColor=C['primary'], fontName='Helvetica-Bold', spaceAfter=3)
    sub_s    = ParagraphStyle('s', fontSize=11, textColor=C['muted'], spaceAfter=2)
    date_s   = ParagraphStyle('d', fontSize=9,  textColor=C['muted'], spaceAfter=10, fontName='Helvetica-Oblique')
    sec_s    = ParagraphStyle('sec', fontSize=13, textColor=C['primary'], fontName='Helvetica-Bold', spaceBefore=10, spaceAfter=5)
    cell_s   = ParagraphStyle('cell', fontSize=7.5, leading=10)
    cell_c   = ParagraphStyle('cellc', fontSize=7.5, leading=10, alignment=TA_CENTER)

    story = []

    # Encabezado
    story.append(Paragraph('Reporte de Entregables', title_s))
    story.append(Paragraph(f"División: <b>{division_filtro or 'Todas las divisiones'}</b>", sub_s))
    story.append(Paragraph(f"Generado: {datetime.now().strftime('%d/%m/%Y  %H:%M')}", date_s))
    story.append(HRFlowable(width='100%', thickness=1.5, color=C['primary'], spaceAfter=8))

    # KPIs
    story.append(Paragraph('Indicadores clave', sec_s))
    kpi_headers = ['Total', 'Verde', 'Amarillo', 'Rojo', 'FRD sin URL', 'Sin firma', 'Sin aceptación', 'FRDs completados']
    kpi_vals    = [str(len(filtered)), str(sem_verde), str(sem_amarillo), str(sem_rojo),
                   str(pendientes_frd), str(pendientes_firma), str(pendientes_aceptacion), f'{frd_pct}%']
    kpi_table = Table([kpi_headers, kpi_vals], colWidths=[3.3*cm]*8)
    kpi_table.setStyle(TableStyle([
        ('BACKGROUND',    (0,0), (-1,0), C['primary']),
        ('TEXTCOLOR',     (0,0), (-1,0), C['white']),
        ('FONTNAME',      (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE',      (0,0), (-1,0), 8),
        ('ALIGN',         (0,0), (-1,-1), 'CENTER'),
        ('VALIGN',        (0,0), (-1,-1), 'MIDDLE'),
        ('FONTNAME',      (0,1), (-1,1), 'Helvetica-Bold'),
        ('FONTSIZE',      (0,1), (-1,1), 15),
        ('TEXTCOLOR',     (1,1), (1,1),  C['success']),
        ('TEXTCOLOR',     (2,1), (2,1),  C['warning']),
        ('TEXTCOLOR',     (3,1), (3,1),  C['danger']),
        ('TEXTCOLOR',     (4,1), (4,1),  C['warning']),
        ('TEXTCOLOR',     (5,1), (5,1),  hex2color('#0284C7')),
        ('TEXTCOLOR',     (6,1), (6,1),  C['danger']),
        ('TEXTCOLOR',     (7,1), (7,1),  C['success']),
        ('BACKGROUND',    (0,1), (-1,1), C['light_bg']),
        ('GRID',          (0,0), (-1,-1), 0.4, hex2color('#D1D5DB')),
        ('TOPPADDING',    (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('ROWHEIGHT',     (0,0), (-1,0), 20),
        ('ROWHEIGHT',     (0,1), (-1,1), 28),
    ]))
    story.append(kpi_table)
    story.append(Spacer(1, 0.5*cm))

    # Tabla entregables
    story.append(Paragraph('Entregables', sec_s))
    ent_hdrs = ['Folio', 'Nombre', 'Estatus', 'Semáforo', 'Avance', 'Responsable', 'F. Entrega']
    ent_data = [ent_hdrs]

    sem_colors_rl = {
        'verde':    C['success'],
        'amarillo': C['warning'],
        'rojo':     C['danger'],
    }

    row_styles = [
        ('BACKGROUND',    (0,0), (-1,0), C['primary']),
        ('TEXTCOLOR',     (0,0), (-1,0), C['white']),
        ('FONTNAME',      (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE',      (0,0), (-1,-1), 7.5),
        ('ALIGN',         (0,0), (-1,-1), 'CENTER'),
        ('ALIGN',         (1,1), (1,-1), 'LEFT'),
        ('ALIGN',         (5,1), (5,-1), 'LEFT'),
        ('VALIGN',        (0,0), (-1,-1), 'MIDDLE'),
        ('GRID',          (0,0), (-1,-1), 0.3, hex2color('#E5E7EB')),
        ('TOPPADDING',    (0,0), (-1,-1), 4),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('ROWHEIGHT',     (0,0), (-1,0), 20),
    ]

    for i, e in enumerate(filtered, 1):
        sem = e.get('semaforo_cliente', '')
        avance = e.get('porcentaje_avance') or 0
        bg = hex2color('#F9FAFB') if i % 2 == 0 else C['white']
        row_styles.append(('BACKGROUND', (0,i), (-1,i), bg))
        if sem in sem_colors_rl:
            row_styles.append(('TEXTCOLOR', (3,i), (3,i), sem_colors_rl[sem]))
        avance_color = C['success'] if avance >= 80 else (C['warning'] if avance >= 40 else C['danger'])
        row_styles.append(('TEXTCOLOR', (4,i), (4,i), avance_color))

        ent_data.append([
            e.get('folio') or '—',
            Paragraph(e.get('nombre') or '', cell_s),
            e.get('estatus') or '',
            _sem_label(sem),
            f"{avance}%",
            Paragraph(e.get('responsable') or '—', cell_s),
            _fmt_fecha(e.get('fecha_entrega')),
        ])

    cw = [1.8*cm, 7.5*cm, 2.6*cm, 2.2*cm, 1.6*cm, 5.5*cm, 2.4*cm]
    ent_table = Table(ent_data, colWidths=cw, repeatRows=1)
    ent_table.setStyle(TableStyle(row_styles))
    story.append(ent_table)

    doc.build(story)
    buf.seek(0)
    return buf


# =============================================================================
#  POWERPOINT
# =============================================================================
def generate_pptx(filtered, division_filtro, sem_verde, sem_amarillo, sem_rojo,
                  pendientes_frd, pendientes_firma, pendientes_aceptacion, frd_pct, total_frd):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb(h):
        h = h.lstrip('#')
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    RC = {k: rgb(v) for k, v in _C.items()}

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    def rect(sl, x, y, w, h, fill=None, line=None, lw=None, radius=False):
        from pptx.util import Inches, Pt
        from pptx.oxml.ns import qn
        shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid() if fill else shp.fill.background()
        if fill:
            shp.fill.fore_color.rgb = fill
        if line:
            shp.line.color.rgb = line
            if lw:
                shp.line.width = Pt(lw)
        else:
            shp.line.fill.background()
        return shp

    def txt(sl, text, x, y, w, h, size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
            italic=False, wrap=True, va='top'):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tb.word_wrap = wrap
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return tb

    div_label = division_filtro or 'Todas las divisiones'

    # ── Slide 1: Portada ─────────────────────────────────────────
    s1 = slide()
    rect(s1, 0, 0, 13.33, 7.5, fill=RC['bg'])
    rect(s1, 0, 0, 0.22, 7.5, fill=RC['primary'])
    rect(s1, 0.22, 3.15, 13.11, 0.04, fill=RC['primary'])

    txt(s1, 'Reporte de Entregables', 0.6, 1.6, 11, 1.1,
        size=38, bold=True, color=RC['primary'], align=PP_ALIGN.LEFT)
    txt(s1, f'División: {div_label}', 0.6, 3.4, 10, 0.55,
        size=18, color=RC['muted'], align=PP_ALIGN.LEFT)
    txt(s1, f"Generado: {datetime.now().strftime('%d de %B de %Y')}", 0.6, 4.0, 10, 0.4,
        size=12, color=RC['muted'], italic=True)

    rect(s1, 10.2, 2.8, 2.7, 1.5, fill=RC['primary'])
    txt(s1, str(len(filtered)), 10.2, 2.85, 2.7, 0.9,
        size=42, bold=True, color=RC['white'], align=PP_ALIGN.CENTER)
    txt(s1, 'entregables', 10.2, 3.75, 2.7, 0.4,
        size=11, color=RC['white'], align=PP_ALIGN.CENTER)

    # ── Slide 2: KPIs actividad + FRD ────────────────────────────
    s2 = slide()
    rect(s2, 0, 0, 13.33, 7.5, fill=RC['bg'])
    rect(s2, 0, 0, 0.22, 7.5, fill=RC['primary'])

    txt(s2, 'Indicadores', 0.55, 0.22, 12, 0.55,
        size=24, bold=True, color=RC['primary'])
    txt(s2, div_label, 0.55, 0.78, 10, 0.35, size=11, color=RC['muted'])

    activity = [
        ('Total',    str(len(filtered)), RC['primary'], 'entregables'),
        ('Verde',    str(sem_verde),     RC['success'],  'a tiempo'),
        ('Amarillo', str(sem_amarillo),  RC['warning'],  'en riesgo'),
        ('Rojo',     str(sem_rojo),      RC['danger'],   'atrasados'),
    ]
    frd_cards = [
        ('FRD sin URL',      str(pendientes_frd),          RC['warning']),
        ('Sin firma',        str(pendientes_firma),         rgb('#0284C7')),
        ('Sin aceptación',   str(pendientes_aceptacion),   RC['danger']),
        ('FRD completados',  f'{frd_pct}%',                RC['success']),
    ]

    x0, cw_card, gap = 0.55, 2.98, 0.12
    for i, (label, val, color, sub) in enumerate(activity):
        x = x0 + i * (cw_card + gap)
        rect(s2, x, 1.3, cw_card, 2.1, fill=color)
        txt(s2, label, x, 1.42, cw_card, 0.38, size=11, color=RC['white'], align=PP_ALIGN.CENTER)
        txt(s2, val, x, 1.8, cw_card, 0.85, size=36, bold=True, color=RC['white'], align=PP_ALIGN.CENTER)
        txt(s2, sub, x, 2.7, cw_card, 0.35, size=9, color=RC['white'], align=PP_ALIGN.CENTER, italic=True)

    for i, (label, val, color) in enumerate(frd_cards):
        x = x0 + i * (cw_card + gap)
        rect(s2, x, 3.65, cw_card, 1.55, fill=rgb('#EEF2FF'))
        rect(s2, x, 3.65, 0.07, 1.55, fill=color)
        txt(s2, label, x+0.18, 3.75, cw_card-0.25, 0.38, size=9, color=RC['muted'])
        txt(s2, val, x+0.18, 4.12, cw_card-0.25, 0.72, size=26, bold=True, color=color)

    # ── Slides 3+: Tabla de entregables (14 filas por slide) ─────
    max_rows = 14
    chunks = [filtered[i:i+max_rows] for i in range(0, max(len(filtered), 1), max_rows)]

    sem_colors_pptx = {'verde': RC['success'], 'amarillo': RC['warning'], 'rojo': RC['danger']}
    hdrs  = ['Folio', 'Nombre', 'Estatus', 'Semáforo', 'Avance', 'Responsable', 'F. Entrega']
    col_w = [1.0,     3.8,      1.6,       1.3,        1.0,      2.5,           1.6]

    for page, chunk in enumerate(chunks):
        s = slide()
        rect(s, 0, 0, 13.33, 7.5, fill=RC['bg'])
        rect(s, 0, 0, 0.22, 7.5, fill=RC['primary'])

        title = 'Entregables' + (f'  ({page+1}/{len(chunks)})' if len(chunks) > 1 else '')
        txt(s, title, 0.55, 0.2, 12, 0.5, size=20, bold=True, color=RC['primary'])
        txt(s, div_label, 0.55, 0.7, 10, 0.3, size=10, color=RC['muted'])

        x0_t = 0.55
        y0_t = 1.1
        h_hdr = 0.4
        h_row = (7.5 - y0_t - h_hdr - 0.1) / max_rows

        # Header
        xc = x0_t
        for h, cw in zip(hdrs, col_w):
            rect(s, xc, y0_t, cw, h_hdr, fill=RC['primary'])
            txt(s, h, xc+0.04, y0_t+0.06, cw-0.08, h_hdr-0.1,
                size=8, bold=True, color=RC['white'], align=PP_ALIGN.CENTER)
            xc += cw

        for ri, e in enumerate(chunk):
            yr = y0_t + h_hdr + ri * h_row
            row_bg = rgb('#F9FAFB') if ri % 2 == 0 else RC['white']
            sem = e.get('semaforo_cliente', '')
            avance = e.get('porcentaje_avance') or 0

            row_vals = [
                e.get('folio') or '—',
                (e.get('nombre') or '')[:45],
                e.get('estatus') or '',
                _sem_label(sem),
                f"{avance}%",
                (e.get('responsable') or '—')[:28],
                _fmt_fecha(e.get('fecha_entrega')),
            ]
            xc = x0_t
            for ci, (val, cw) in enumerate(zip(row_vals, col_w)):
                rect(s, xc, yr, cw, h_row, fill=row_bg,
                     line=rgb('#E5E7EB'), lw=0.3)
                color = RC['dark']
                if ci == 3 and sem in sem_colors_pptx:
                    color = sem_colors_pptx[sem]
                elif ci == 4:
                    color = RC['success'] if avance >= 80 else (RC['warning'] if avance >= 40 else RC['danger'])
                align = PP_ALIGN.LEFT if ci in (1, 5) else PP_ALIGN.CENTER
                txt(s, str(val), xc+0.05, yr+0.03, cw-0.08, h_row-0.04,
                    size=7.5, color=color, align=align)
                xc += cw

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
